CREATE DATABASE IF NOT EXISTS HRDATA_CORTEX_SEARCH;

USE DATABASE HRDATA_CORTEX_SEARCH;
CREATE SCHEMA IF NOT EXISTS PUBLIC;

-- CREATE OR REPLACE STAGE HRDATA_CORTEX_SEARCH.PUBLIC.RAW
--     DIRECTORY = (ENABLE = TRUE)
--     ENCRYPTION = (TYPE = 'SNOWFLAKE_SSE');

create or replace function chunking(file_url string , relative_url string)
returns table (chunk string) 
language python
runtime_version = '3.10'
handler = 'file_text_chunker'
packages = ('snowflake-snowpark-python','langchain', 'python-pptx','python-docx','pypdf2')
as
$$
from snowflake.snowpark.types import StringType, StructField, StructType
from langchain.text_splitter import RecursiveCharacterTextSplitter
from snowflake.snowpark.files import SnowflakeFile
import io
from pptx import Presentation
import logging
import pandas as pd
import PyPDF2
from docx import Document
import os

class file_text_chunker:    
     def read_pdf(self, file_url: str) -> str:
        logger = logging.getLogger("udf_logger")
        logger.info(f"Opening file {file_url}")
    
        with SnowflakeFile.open(file_url, 'rb') as f:
            buffer = io.BytesIO(f.readall())
            
        reader = PyPDF2.PdfReader(buffer)   
        text = ""
        for page in reader.pages:
            try:
                text += page.extract_text().replace('\n', ' ').replace('\0', ' ')
            except:
                text = "Unable to Extract"
                logger.warn(f"Unable to extract from file {file_url}, page {page}")
        
        return text

     def read_docx(self, file_url: str) -> str:
        logger = logging.getLogger("udf_logger")
        logger.info(f"Opening file {file_url}")
        with SnowflakeFile.open(file_url, 'rb') as f:
            buffer = io.BytesIO(f.read())
        document = Document(buffer)
        text = ""
        for paragraph in document.paragraphs:
            text += paragraph.text.replace('\n', ' ').replace('\0', ' ') + " "
        
        return text
        
     def read_pptx(self, file_url: str) -> str:
        logger = logging.getLogger("udf_logger")
        logger.info(f"Opening file {file_url}")
        
        with SnowflakeFile.open(file_url, 'rb') as f:
            buffer = io.BytesIO(f.readall())
        
        presentation = Presentation(buffer)
        text = ""
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text.replace('\n','').replace('\0', ' ') + " "

        return text
        
     def process(self, file_url, relative_url):
        if relative_url.split(".")[-1] == "pptx":
            text = self.read_pptx(file_url)  
        elif relative_url.split(".")[-1] == "docx":
            text = self.read_docx(file_url)  
        elif relative_url.split(".")[-1] == "pdf":
            text = self.read_pdf(file_url)  

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=4000,  # Adjust this as you see fit
            chunk_overlap=400,  # This lets text have some form of overlap. Useful for keeping chunks contextual
            length_function=len
        )
    
        chunks = text_splitter.split_text(text)
        df = pd.DataFrame(chunks, columns=['chunks'])
        
        yield from df.itertuples(index=False, name=None)
$$;


CREATE OR REPLACE TABLE HRDATA_CORTEX_SEARCH.PUBLIC.docs_chunks_table AS
    SELECT
        relative_path,
        build_scoped_file_url(@HRDATA_CORTEX_SEARCH.PUBLIC.RAW, relative_path) AS file_url,
        -- preserve file title information by concatenating relative_path with the chunk
        CONCAT(relative_path, ': ', func.chunk) AS chunk,
        'English' AS language
    FROM
        directory(@HRDATA_CORTEX_SEARCH.PUBLIC.RAW),
        TABLE(HRDATA_CORTEX_SEARCH.PUBLIC.CHUNKING(build_scoped_file_url(@HRDATA_CORTEX_SEARCH.PUBLIC.RAW, relative_path),relative_path)) AS func;


CREATE OR REPLACE CORTEX SEARCH SERVICE HRDATA_CORTEX_SEARCH.PUBLIC.RAW_INDEX
    ON chunk
    ATTRIBUTES language
    WAREHOUSE = compute_wh
    TARGET_LAG = '1 hour'
    AS (
    SELECT
        chunk,
        relative_path,
        file_url,
        language
    FROM HRDATA_CORTEX_SEARCH.PUBLIC.docs_chunks_table
    );

